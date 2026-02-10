import os
from typing import List, Dict
import PyPDF2
import docx
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup
import json
import csv


class DocumentProcessor:
    """Extracts text from various document formats"""
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._extract_pdf,
            '.txt': self._extract_txt,
            '.docx': self._extract_docx,
            '.doc': self._extract_docx,
            '.xlsx': self._extract_xlsx,
            '.xls': self._extract_xlsx,
            '.pptx': self._extract_pptx,
            '.ppt': self._extract_pptx,
            '.md': self._extract_txt,
            '.html': self._extract_html,
            '.htm': self._extract_html,
            '.csv': self._extract_csv,
            '.json': self._extract_json,
            '.xml': self._extract_xml,
            '.rtf': self._extract_txt,
        }
    
    def process_document(self, file_path: str) -> Dict[str, any]:
        """
        Process a document and extract its text content
        
        Args:
            file_path: Path to the document
            
        Returns:
            Dictionary with document metadata and content
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        try:
            text_content = self.supported_formats[file_ext](file_path)
            
            return {
                'file_path': file_path,
                'file_name': os.path.basename(file_path),
                'file_type': file_ext,
                'content': text_content,
                'size': os.path.getsize(file_path),
                'success': True,
                'error': None
            }
        except Exception as e:
            return {
                'file_path': file_path,
                'file_name': os.path.basename(file_path),
                'file_type': file_ext,
                'content': '',
                'size': os.path.getsize(file_path),
                'success': False,
                'error': str(e)
            }
    
    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF"""
        text = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text.append(f"--- Page {page_num + 1} ---\n{page_text}")
        return "\n\n".join(text)
    
    def _extract_txt(self, file_path: str) -> str:
        """Extract text from plain text files"""
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        raise ValueError(f"Unable to decode file with supported encodings")
    
    def _extract_docx(self, file_path: str) -> str:
        """Extract text from Word documents"""
        doc = docx.Document(file_path)
        text = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                text.append(" | ".join(row_text))
        
        return "\n\n".join(text)
    
    def _extract_xlsx(self, file_path: str) -> str:
        """Extract text from Excel files"""
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text.append(f"=== Sheet: {sheet_name} ===")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else '' for cell in row]
                if any(row_text):  # Skip empty rows
                    text.append(" | ".join(row_text))
        
        return "\n\n".join(text)
    
    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentations"""
        prs = Presentation(file_path)
        text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text.append(f"=== Slide {slide_num} ===")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        
        return "\n\n".join(text)
    
    def _extract_html(self, file_path: str) -> str:
        """Extract text from HTML files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            soup = BeautifulSoup(file.read(), 'lxml')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return text
    
    def _extract_csv(self, file_path: str) -> str:
        """Extract text from CSV files"""
        text = []
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            csv_reader = csv.reader(file)
            for row in csv_reader:
                text.append(" | ".join(row))
        return "\n".join(text)
    
    def _extract_json(self, file_path: str) -> str:
        """Extract text from JSON files"""
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
            return json.dumps(data, indent=2)
    
    def _extract_xml(self, file_path: str) -> str:
        """Extract text from XML files"""
        with open(file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file.read(), 'lxml')
            return soup.get_text()
    
    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """
        Split text into overlapping chunks
        
        Args:
            text: Text to chunk
            chunk_size: Size of each chunk in characters
            overlap: Number of overlapping characters between chunks
            
        Returns:
            List of text chunks
        """
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence endings
                for sep in ['. ', '.\n', '! ', '!\n', '? ', '?\n']:
                    last_sep = text[start:end].rfind(sep)
                    if last_sep != -1:
                        end = start + last_sep + len(sep)
                        break
            
            chunks.append(text[start:end].strip())
            start = end - overlap
        
        return chunks
