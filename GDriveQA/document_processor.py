"""Document processing for various file formats."""
from pathlib import Path
from typing import List, Dict, Optional
import PyPDF2
import docx
from pptx import Presentation
from openpyxl import load_workbook
from bs4 import BeautifulSoup
import markdown


class DocumentProcessor:
    """Process documents and extract text content."""
    
    @staticmethod
    def extract_text(file_path: Path) -> Optional[str]:
        """Extract text from a file based on its extension."""
        suffix = file_path.suffix.lower()
        
        extractors = {
            '.pdf': DocumentProcessor._extract_pdf,
            '.docx': DocumentProcessor._extract_docx,
            '.pptx': DocumentProcessor._extract_pptx,
            '.xlsx': DocumentProcessor._extract_xlsx,
            '.txt': DocumentProcessor._extract_txt,
            '.md': DocumentProcessor._extract_markdown,
            '.html': DocumentProcessor._extract_html,
            '.rtf': DocumentProcessor._extract_txt,
        }
        
        extractor = extractors.get(suffix)
        if not extractor:
            return None
        
        try:
            return extractor(file_path)
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return None
    
    @staticmethod
    def _extract_pdf(file_path: Path) -> str:
        """Extract text from PDF."""
        text = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text.append(page.extract_text())
        return '\n\n'.join(text)
    
    @staticmethod
    def _extract_docx(file_path: Path) -> str:
        """Extract text from DOCX."""
        doc = docx.Document(file_path)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return '\n\n'.join(text)
    
    @staticmethod
    def _extract_pptx(file_path: Path) -> str:
        """Extract text from PPTX."""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n\n'.join(text)
    
    @staticmethod
    def _extract_xlsx(file_path: Path) -> str:
        """Extract text from XLSX."""
        wb = load_workbook(file_path, read_only=True)
        text = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text.append(f"=== Sheet: {sheet_name} ===")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip():
                    text.append(row_text)
        
        return '\n'.join(text)
    
    @staticmethod
    def _extract_txt(file_path: Path) -> str:
        """Extract text from plain text file."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            return file.read()
    
    @staticmethod
    def _extract_markdown(file_path: Path) -> str:
        """Extract text from Markdown."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            md_content = file.read()
            # Convert markdown to HTML then extract text
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, 'html.parser')
            return soup.get_text()
    
    @staticmethod
    def _extract_html(file_path: Path) -> str:
        """Extract text from HTML."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            soup = BeautifulSoup(file.read(), 'html.parser')
            return soup.get_text()
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks."""
        if not text:
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = start + chunk_size
            
            # Try to break at sentence boundary
            if end < text_length:
                # Look for sentence endings
                for punct in ['. ', '.\n', '! ', '!\n', '? ', '?\n']:
                    last_punct = text.rfind(punct, start, end)
                    if last_punct != -1:
                        end = last_punct + 1
                        break
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - overlap if end < text_length else text_length
        
        return chunks
    
    @staticmethod
    def process_document(file_path: Path, chunk_size: int = 1000, 
                        overlap: int = 200) -> Dict:
        """Process a document: extract text and create chunks."""
        text = DocumentProcessor.extract_text(file_path)
        
        if not text:
            return {
                'file_path': str(file_path),
                'file_name': file_path.name,
                'success': False,
                'text': None,
                'chunks': []
            }
        
        chunks = DocumentProcessor.chunk_text(text, chunk_size, overlap)
        
        return {
            'file_path': str(file_path),
            'file_name': file_path.name,
            'success': True,
            'text': text,
            'chunks': chunks,
            'num_chunks': len(chunks)
        }
